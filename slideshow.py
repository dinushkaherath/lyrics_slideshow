# slideshow.py
from typing import List, Tuple, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide

from utils import calculate_dynamic_font_size, estimate_wrapped_lines

class LyricsSlideshow:
    def __init__(self, config: Dict[str, Any]) -> None:
        """
        Initialize PowerPoint presentation with slide size and blank layout.
        
        Args:
            config: A configuration dictionary from config.py
        """
        self.config = config
        self.prs = Presentation()

        # Set slide size from config
        self.prs.slide_width = self.config["POSITION"]["slide_width"]
        self.prs.slide_height = self.config["POSITION"]["slide_height"]
        self.blank_layout = self.prs.slide_layouts[6]

    def _add_icon(self, slide: Slide, target_slide: Optional[Slide], icon_path: str, icon_type: str) -> None:
        """Add a clickable icon to the slide."""
        width, height = self.config["ICON_SIZES"][icon_type]
        left, top = self.config["ICON_POSITIONS"][icon_type]
        try:
            pic = slide.shapes.add_picture(icon_path, left, top, width=width, height=height)
            if target_slide:
                pic.click_action.target_slide = target_slide
        except Exception as e:
            print(f"⚠️ Error adding {icon_type} icon ({icon_path}): {e}")

    def _add_header_background(self, slide: Slide) -> object:
        """Adds a colored header background strip."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Length(0),  # Left
            self.config["POSITION"]["bottom_margin"],
            self.config["POSITION"]["slide_width"],
            self.config["POSITION"]["header_height"] + Inches(0.1)
        )
        fill = shape.fill
        if self.config["backheadfil"] == "y":
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.config["COLOR"]["header_bg"]
        else:
            shape.fill.background()
        shape.line.fill.background()  # No border
        return shape

    def _add_text_box(self,
                      slide: Slide,
                      text: str,
                      left: Length,
                      top: Length,
                      width: Length,
                      height: Length,
                      font_size: Optional[Pt] = None,
                      horizontal_alignment=PP_ALIGN.CENTER,
                      vertical_alignment=MSO_VERTICAL_ANCHOR.MIDDLE,
                      font_type: str = "body",
                      is_chorus: bool = False,
                      is_header: bool = False) -> object:
        """Add a formatted text box with optional dynamic sizing."""
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Let text wrap within the box
        text_frame.vertical_anchor = vertical_alignment

        # Split text into lines and add each line as a separate paragraph
        lines = text.split('\n')

        # Use dynamic font size if none is provided
        if font_size is None:
            font_size = calculate_dynamic_font_size(text, is_chorus)

        # Helper to apply font styles
        def style_font(run, italic=False, header=False):
            run.font.size = font_size
            run.font.name = self.config["FONT"]["header"] if header else self.config["FONT"][font_type]
            run.font.color.rgb = self.config["COLOR"]["header_text"] if header else self.config["COLOR"]["text"]
            run.font.italic = italic

        # First paragraph
        first_para = text_frame.paragraphs[0]
        first_para.text = lines[0].upper() if is_header else lines[0]
        first_para.alignment = horizontal_alignment
        first_para.line_spacing = 1.0
        
        if first_para.runs:
             style_font(first_para.runs[0], italic=is_chorus, header=is_header)

        # Remaining paragraphs
        for line in lines[1:]:
            para = text_frame.add_paragraph()
            para.text = line.upper() if is_header else line
            para.alignment = horizontal_alignment
            para.line_spacing = 1.0
            if para.runs:
                style_font(para.runs[0], italic=is_chorus, header=is_header)
        
        return shape

    def _add_title_slide(self) -> Slide:
        """Creates and adds the main title slide."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.config["COLOR"]["bg"]
        
        self._add_text_box(
            slide=slide,
            text="Song Lyrics Slideshow",
            left=self.config["POSITION"]["left_margin"],
            top=self.config["POSITION"]["title_top"],
            width=self.config["POSITION"]["slide_width"],
            height=self.config["POSITION"]["title_height"],
            font_size=self.config["SIZE"]["title"],
            font_type="title"
        )
        return slide

    def _add_song_list_slide(self,
                             song_titles: List[str],
                             song_slide_map: Dict[int, Slide]
                             ) -> Tuple[Slide, Dict[int, int]]:
        """Creates a slide listing all songs with clickable links."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.config["COLOR"]["header_bg"]

        # Grid configuration
        grid = self.config["GRID"]
        num_columns = grid["columns"]
        box_width = grid["box_width"]
        box_height = grid["box_height"]
        
        number_index_map = {}

        for index, title in enumerate(song_titles):
            col = index % num_columns
            row = index // num_columns
            song_number = index + 1
            num_str = f"{song_number:2}"  # Aligns 1-9 and 10+
            full_title = f"{num_str} – {title}"

            left = grid["start_left"] + col * (box_width + grid["spacing_x"])
            top = grid["start_top"] + row * (box_height + grid["spacing_y"])

            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.config["COLOR"]["header_bg"]
            shape.line.color.rgb = self.config["COLOR"]["text"]
            shape.line.width = Pt(0.75)

            text_frame = shape.text_frame
            text_frame.text = full_title
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = self.config["SIZE"]["song_list"]
            run.font.name = self.config["FONT"]["body"]
            run.font.color.rgb = self.config["COLOR"]["text"]

            try:
                shape.click_action.target_slide = song_slide_map[index]
            except Exception as e:
                print(f"Error linking '{full_title}' to slide: {e}")

            number_index_map[song_number] = index

        return slide, number_index_map

    def _add_alpha_index_slide(self,
                               alpha_order: List[Tuple[int, str]],
                               song_slide_map: Dict[int, Slide],
                               number_index_map: Dict[int, int]
                               ) -> Slide:
        """Creates an alphabetically ordered index slide."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.config["COLOR"]["header_bg"]

        grid = self.config["GRID"]
        num_columns = grid["columns"]
        box_width = grid["box_width"]
        box_height = grid["box_height"]

        for index, (song_number, title) in enumerate(alpha_order):
            col = index % num_columns
            row = index // num_columns
            num_str = f"{song_number:2}"
            full_title = f"{num_str} – {title}"

            left = grid["start_left"] + col * (grid["box_width"] + grid["spacing_x"])
            top = grid["start_top"] + row * (grid["box_height"] + grid["spacing_y"])

            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.config["COLOR"]["header_bg"]
            shape.line.color.rgb = self.config["COLOR"]["text"]
            shape.line.width = Pt(0.75)

            text_frame = shape.text_frame
            text_frame.text = full_title
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = self.config["SIZE"]["song_list"]
            run.font.name = self.config["FONT"]["body"]
            run.font.color.rgb = self.config["COLOR"]["text"]

            if song_number in number_index_map:
                song_index = number_index_map[song_number]
                try:
                    shape.click_action.target_slide = song_slide_map[song_index]
                except Exception as e:
                    print(f"Error linking '{full_title}' in alpha index: {e}")
            else:
                print(f"Warning: No match for song number {song_number} ({title})")

        return slide

    def create_presentation_from_parsed_sections(self,
                                                 songs: List[Tuple[int, str, int, List[Dict]]],
                                                 alpha_order: List[Tuple[int, str]],
                                                 output_file: str = "lyrics_slideshow.pptx"
                                                 ) -> str:
        """
        Create the full presentation with a title slide, song slides, and index slides.
        """
        # Title slide
        self._add_title_slide()

        song_titles = [title for _, title, _, _ in songs]
        song_slide_map = {}  # Map song index -> first slide

        # --- Create Lyric Slides ---
        for index, (song_number, title, chorus_count, sections) in enumerate(songs):
            for slide_index, section in enumerate(sections):
                lines = section["content"].splitlines()
                # Split long sections into chunks of 9 lines (10 was splitting 10, not max 9)
                chunks = [lines[i:i+9] for i in range(0, len(lines), 9)]
                if not chunks: chunks = [[""]] # Handle empty sections

                for chunk_index, chunk in enumerate(chunks):
                    slide = self.prs.slides.add_slide(self.blank_layout)
                    if slide_index == 0 and chunk_index == 0:
                        song_slide_map[index] = slide

                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = self.config["COLOR"]["bg"]
                    self._add_header_background(slide)

                    # Header left: Song number + title
                    self._add_text_box(
                        slide=slide, text=f"{song_number}: {title}",
                        left=self.config["POSITION"]["left_margin"], top=self.config["POSITION"]["bottom_margin"],
                        width=self.config["POSITION"]["header_width_left"], height=self.config["POSITION"]["header_height"],
                        font_size=self.config["SIZE"]["header"], horizontal_alignment=PP_ALIGN.LEFT,
                        is_header=True, font_type="header"
                    )

                    # Header right: section label
                    section_type = section["type"].upper()
                    if section_type == "CHORUS":
                        section_label = f"CHORUS {section['number']}" if chorus_count > 1 else "CHORUS"
                    else:
                        section_label = f"STANZA {section['number']}"
                    if len(chunks) > 1:
                        section_label += f" ({chunk_index + 1}/{len(chunks)})"

                    self._add_text_box(
                        slide=slide, text=section_label,
                        left=self.config["POSITION"]["header_left"], top=self.config["POSITION"]["bottom_margin"],
                        width=self.config["POSITION"]["header_width_right"], height=self.config["POSITION"]["header_height"],
                        font_size=self.config["SIZE"]["header"], horizontal_alignment=PP_ALIGN.RIGHT,
                        is_header=True, font_type="header"
                    )

                    # Add home icon (target TBD)
                    self._add_icon(slide, target_slide=None, icon_path="assets/home.png", icon_type="home")

                    # Add lyrics box
                    is_chorus = (section["type"] == "chorus")
                    self._add_text_box(
                        slide=slide, text="\n".join(chunk),
                        left=self.config["POSITION"]["left_margin"], top=self.config["POSITION"]["lyrics_top"],
                        width=self.config["POSITION"]["lyrics_width"], height=self.config["POSITION"]["lyrics_height"],
                        font_size=self.config["SIZE"]["stanza"] if not is_chorus else self.config["SIZE"]["chorus"],
                        vertical_alignment=MSO_VERTICAL_ANCHOR.TOP,
                        is_chorus=is_chorus
                    )

                    # Add restart icon on last slide of song
                    if slide_index == len(sections) - 1 and chunk_index == len(chunks) - 1:
                         self._add_icon(slide, target_slide=song_slide_map[index], icon_path="assets/restart.png", icon_type="restart")


        # --- Create Index Slides ---
        song_list_slide, title_index_map = self._add_song_list_slide(song_titles, song_slide_map)
        alpha_index_slide = self._add_alpha_index_slide(alpha_order, song_slide_map, title_index_map)

        # --- Patch Home Icons ---
        # Now that song_list_slide exists, link all home icons to it
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "click_action"):
                    # Find shapes with a click action but no target (our home icons)
                    if getattr(shape.click_action, "target_slide", None) is None:
                        shape.click_action.target_slide = song_list_slide

        self.prs.save(output_file)
        return output_file