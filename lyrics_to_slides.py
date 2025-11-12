from typing import List, Tuple, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.slide import Slide

def configure_defaults():
    """
    Presents a menu to edit default presentation settings interactively.
    The user can review all options, modify specific ones, and confirm when done.
    """

    # Default values
    FONT = {
        "title": "Grathik",
        "body": "Grathik",
        "header": "Grathik",
    }

    SIZE = {
        "title": 44,
        "header": 24,
        "stanza": 44,
        "chorus": 44,
        "song_list": 14,
    }

    COLOR = {
        "bg": "2D1412",
        "header_bg": "BE4736",
        "text": "FFFFFF",
        "header_text": "FFFFFF",
    }

    OPTIONS = {
        1: ("Fonts", FONT),
        2: ("Font Sizes", SIZE),
        3: ("Colors", COLOR),
        4: ("Header Background Strip", {"enabled": "y"}),
    }

    def print_options():
        print("\n=== Configuration Menu ===")
        for num, (label, values) in OPTIONS.items():
            print(f"\n[{num}] {label}:")
            for k, v in values.items():
                print(f"   {k}: {v}")
        print("\nPress ENTER to confirm and continue.\n")

    # --- Interactive Loop --- #
    while True:
        print_options()
        choice = input("Enter the number of the section to edit (or ENTER to finish): ").strip()

        if not choice:
            break

        if not choice.isdigit() or int(choice) not in OPTIONS:
            print("Invalid choice. Please try again.")
            continue

        section_name, section_dict = OPTIONS[int(choice)]
        print(f"\nEditing {section_name}:")

        for key, value in section_dict.items():
            new_val = input(f"  {key} [{value}]: ").strip()
            if new_val:
                section_dict[key] = new_val

    # Convert color strings to RGBColor
    COLOR = {k: RGBColor.from_string(v) for k, v in COLOR.items()}

    # Rebuild SIZE with Pt
    from pptx.util import Pt
    SIZE = {k: Pt(v) for k, v in SIZE.items()}

    backheadfil = OPTIONS[4][1]["enabled"].lower()
    return FONT, SIZE, COLOR, backheadfil

FONT, SIZE, COLOR, backheadfil = configure_defaults()

POSITION = {
    # Slide dimensions
    "slide_width": Inches(13.333),
    "slide_height": Inches(7.5),

    # Margins
    "left_margin": Inches(0.3),
    "right_margin": Inches(9.7),
    "bottom_margin": Inches(7),

    # Lyrics text box
    "lyrics_top": Inches(0),
    "lyrics_height": Inches(5),
    "lyrics_width": Inches(12.7),

    # Header
    "header_height": Inches(0.4),
    "header_width_left": Inches(10.7),
    "header_width_right": Inches(3),
    "header_left": Inches(9.5),
    "header_top": Inches(7),
    "title_top": Inches(3),
    "title_height": Inches(2),
}

ICON_SIZES = {
    "home": (Inches(0.4), Inches(0.4)),
    "restart": (Inches(0.93), Inches(0.75)),
}

ICON_POSITIONS = {
    "home": (
        POSITION["slide_width"] - ICON_SIZES["home"][0] - Inches(0.2),
        POSITION["bottom_margin"]
    ),
    "restart": (Inches(12.06), Inches(6)),
}

GRID = {
    # Song list layout
    "columns": 3,
    "box_width": Inches(4.45),
    "box_height": Inches(0.5),
    "spacing_x": Inches(0),
    "spacing_y": Inches(0),
    "start_left": Inches(0),
    "start_top": Inches(0),
}

class LyricsSlideshow:
    def __init__(self) -> None:
        """Initialize PowerPoint presentation with slide size and blank layout."""
        self.prs = Presentation()

        # Set slide size to 16:9
        self.prs.slide_width = POSITION["slide_width"]
        self.prs.slide_height = POSITION["slide_height"]
        self.blank_layout = self.prs.slide_layouts[6]

    def _calculate_dynamic_font_size(self, text: str, is_chorus: bool = False) -> Pt:
        """
        Estimate font size based on text length to avoid overflow.
        """
        base_size = 28
        min_size = 18
        max_lines = 12 if is_chorus else 14
        max_chars_per_line = 50

        # Count actual visual lines based on wrap
        raw_lines = text.split('\n')
        estimated_lines = sum((len(line) // max_chars_per_line + 1) for line in raw_lines)

        scale = min(1.0, max_lines / max(estimated_lines, 1))
        size = int(base_size * scale)
        return Pt(max(size, min_size))

    def _add_icon(self, slide, target_slide, icon_path: str, icon_type: str) -> None:
        """
        Add a clickable icon to the slide.
        """
        width, height = ICON_SIZES[icon_type]
        left, top = ICON_POSITIONS[icon_type]
        try:
            pic = slide.shapes.add_picture(icon_path, left, top, width=width, height=height)
            if target_slide:
                pic.click_action.target_slide = target_slide
        except Exception as e:
            print(f"⚠️ Error adding {icon_type} icon: {e}")

    def _add_header_background(self, slide) -> object:
        """
        Adds a colored header background strip.
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,  # Left
            POSITION["bottom_margin"],
            POSITION["slide_width"],
            POSITION["header_height"] + Inches(0.1)
        )
        fill = shape.fill
        if backheadfil == "y":
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR["header_bg"]
        else:
            shape.fill.background()
        shape.line.fill.background()  # No border

        return shape

    def _add_text_box(self,
                      slide,
                      text: str,
                      left: float,
                      top: float,
                      width: float,
                      height: float,
                      font_size: Optional[Pt] = None,
                      horizontal_alignment=PP_ALIGN.CENTER,
                      vertical_alignment=MSO_VERTICAL_ANCHOR.MIDDLE,
                      font_type: str = FONT["body"],
                      is_chorus: bool = False,
                      is_header: bool = False) -> object:
        """
        Add a formatted text box with optional dynamic sizing.
        """
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Let text wrap within the box
        text_frame.vertical_anchor = vertical_alignment

        # Split text into lines and add each line as a separate paragraph
        lines = text.split('\n')

        # Compute font size dynamically unless header forces it
        if font_size is None:
            font_size = self._calculate_dynamic_font_size(text, is_chorus)

        # Helper to apply font styles
        def style_font(run, italic=False, header=False):
            run.font.size = font_size
            run.font.name = FONT["header"] if header else font_type
            run.font.color.rgb = COLOR["header_text"] if header else COLOR["text"]
            run.font.italic = italic

        # First paragraph
        first_para = text_frame.paragraphs[0]
        first_para.text = lines[0].upper() if is_header else lines[0]
        first_para.alignment = horizontal_alignment
        first_para.line_spacing = 1.0
        style_font(first_para.runs[0], italic=is_chorus, header=is_header)

        # Remaining paragraphs
        for line in lines[1:]:
            para = text_frame.add_paragraph()
            para.text = line.upper() if is_header else line
            para.alignment = horizontal_alignment
            para.line_spacing = 1.0
            style_font(para.runs[0], italic=is_chorus, header=is_header)

        return shape

    def _estimate_wrapped_lines(self, text: str, box_width_inches: float, font_size_pt: float) -> int:
        """
        bruh
        """
        char_width_factor = 0.55  # Empirical average character width factor
        avg_char_width_inches = (font_size_pt / 72.0) * char_width_factor # in points, where 1pt = 1/72 inch

        if avg_char_width_inches == 0:
            return 1

        chars_per_line = box_width_inches / avg_char_width_inches
        lines = text.split('\n')

        estimated_lines = 0
        for line in lines:
            est_line_count = max(1, int(len(line) / chars_per_line) + 1)
            estimated_lines += est_line_count

        return estimated_lines

    def create_presentation_from_parsed_sections(
        self,
        songs: List[Tuple[int, str, int, List[Dict]]],
        alpha_order: List[Tuple[int, str]],
        output_file: str = "lyrics_slideshow.pptx"
    ) -> str:
        """
        Create the full presentation with a title slide, song slides, and a song list slide.
        """
        # Title slide
        title_slide = self.prs.slides.add_slide(self.blank_layout)
        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = COLOR["bg"]
        self._add_text_box(
            slide=title_slide,
            text="Song Lyrics Slideshow",
            left=POSITION["left_margin"],
            top=POSITION["title_top"],
            width=POSITION["slide_width"],
            height=POSITION["title_height"],
            font_size=SIZE["title"],
            font_type=FONT["title"]
        )

        song_titles = [title for _, title, _, _ in songs]

        # Map song index → first slide for restart icons
        song_slide_map = {}

        for index, (song_number, title, chorus_count, sections) in enumerate(songs):
            for slide_index, section in enumerate(sections):
                lines = section["content"].splitlines()
                chunks = [lines[i:i+9] for i in range(0, len(lines), 10)]  # max 9 lines per slide

                for chunk_index, chunk in enumerate(chunks):
                    slide = self.prs.slides.add_slide(self.blank_layout)
                    if slide_index == 0 and chunk_index == 0:
                        song_slide_map[index] = slide

                    # Background & header
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = COLOR["bg"]
                    self._add_header_background(slide)

                    # Header left: Song number + title
                    self._add_text_box(
                        slide=slide,
                        text=f"{song_number}: {title}",
                        left=POSITION["left_margin"],
                        top=POSITION["bottom_margin"],
                        width=POSITION["header_width_left"],
                        height=POSITION["header_height"],
                        font_size=SIZE["header"],
                        horizontal_alignment=PP_ALIGN.LEFT,
                        is_header=True,
                        font_type=FONT["header"]
                    )

                    # Header right: section label
                    section_type = section["type"].upper()
                    if section_type == "CHORUS":
                        section_label = f"CHORUS {section['number']}" if chorus_count > 1 else "CHORUS"
                    else:
                        section_label = f"STANZA {section['number']}"

                    if len(chunks) > 1:
                        section_label += f" ({chunk_index + 1}/{len(chunks)})"

                    self._add_text_box(
                        slide=slide,
                        text=section_label,
                        left=POSITION["header_left"],
                        top=POSITION["bottom_margin"],
                        width=POSITION["header_width_right"],
                        height=POSITION["header_height"],
                        font_size=SIZE["header"],
                        horizontal_alignment=PP_ALIGN.RIGHT,
                        is_header=True,
                        font_type=FONT["header"]
                    )

                    self._add_icon(slide, target_slide=None, icon_path="assets/home.png", icon_type="home")

                    # Add lyrics box (chunk joined back to string)
                    self._add_text_box(
                        slide=slide,
                        text="\n".join(chunk),
                        left=POSITION["left_margin"],
                        top=POSITION["lyrics_top"],
                        width=POSITION["lyrics_width"],
                        height=POSITION["lyrics_height"],
                        font_size=SIZE["stanza"] if section["type"] == "stanza" else SIZE["chorus"],
                        vertical_alignment=MSO_VERTICAL_ANCHOR.TOP,
                        is_chorus=(section["type"] == "chorus")
                    )

                # Debug: Print slides with more than 9 lines
                font_size = SIZE["stanza"] if section['type'] == 'stanza' else SIZE["chorus"]
                wrapped_line_count = self._estimate_wrapped_lines(section['content'], POSITION["lyrics_width"], font_size.pt)

                if wrapped_line_count > 9:
                    print(f"[⚠️ Long Slide] Song {song_number}, slide {len(self.prs.slides)} has {wrapped_line_count} wrapped lines")

                # Add restart icon on last slide of song, linking to first slide of that song
                if slide_index == len(sections) - 1:
                    self._add_icon(slide, target_slide=song_slide_map[index], icon_path="assets/restart.png", icon_type="restart")

        # Song list slide (normal numeric order)
        song_list_slide, title_index_map = self._add_song_list_slide(song_titles, song_slide_map)

        # Alphabetical index slide (added right after)
        alpha_index_slide = self._add_alpha_index_slide(alpha_order, song_slide_map, title_index_map)

        # Patch home icons to link to song list slide
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "click_action"):
                    if getattr(shape.click_action, "target_slide", None) is None:
                        shape.click_action.target_slide = song_list_slide

        self.prs.save(output_file)
        return output_file

    def _add_song_list_slide(self, song_titles: List[str], song_slide_map: Dict[int, Slide]) -> Tuple[Slide, Dict[int, int]]:
        """
        Creates a slide listing all songs with clickable boxes linking to each song's first slide.
        Returns a mapping of song_number -> slide index for later reference.
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Set background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR["header_bg"]

        # Grid configuration
        num_columns = GRID["columns"]
        box_width = GRID["box_width"]
        box_height = GRID["box_height"]
        spacing_x = GRID["spacing_x"]
        spacing_y = GRID["spacing_y"]
        left_start = GRID["start_left"]
        top_start = GRID["start_top"]

        number_index_map = {}

        for index, title in enumerate(song_titles):
            col = index % num_columns
            row = index // num_columns
            song_number = index + 1

            # Add a leading space for single-digit numbers
            num_str = f"{song_number:2}"  # Ensures alignment for 1–9, 10+
            full_title = f"{num_str} – {title}"

            left = left_start + col * (box_width + spacing_x)
            top = top_start + row * (box_height + spacing_y)

            # Create the rectangle (serves as background + border + text container)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
            # Set fill color
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR["header_bg"]
            # Set border (white thin line)
            shape.line.color.rgb = COLOR["header_text"]
            shape.line.width = Pt(0.75)

            # Text styling
            text_frame = shape.text_frame
            text_frame.text = full_title
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.word_wrap = True

            # Format paragraph
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(14)
            run.font.name = FONT["body"]
            run.font.color.rgb = COLOR["text"]

            # Link to the correct song’s first slide
            try:
                shape.click_action.target_slide = song_slide_map[index]
            except Exception as e:
                print(f"Error linking '{full_title}' to slide: {e}")

            # Store mapping of song_number → index
            number_index_map[song_number] = index

        return slide, number_index_map

    def _add_alpha_index_slide(
        self,
        alpha_order: List[Tuple[int, str]],
        song_slide_map: Dict[int, Slide],
        number_index_map: Dict[int, int]
    ) -> Slide:
        """
        Creates an alphabetically ordered index slide with clickable links to each song.
        Uses song_number-based matching for linking.
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        # Background styling
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR["header_bg"]

        # Grid setup
        num_columns = GRID["columns"]
        box_width = GRID["box_width"]
        box_height = GRID["box_height"]
        spacing_x = GRID["spacing_x"]
        spacing_y = GRID["spacing_y"]
        left_start = GRID["start_left"]
        top_start = GRID["start_top"]

        for index, (song_number, title) in enumerate(alpha_order):
            col = index % num_columns
            row = index // num_columns
            
            # Add a leading space for single-digit numbers
            num_str = f"{song_number:2}"  # width of 2 ensures alignment for 1–9, 10+
            full_title = f"{num_str} – {title}"

            left = left_start + col * (box_width + spacing_x)
            top = top_start + row * (box_height + spacing_y)

            # Create text box background
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR["header_bg"]
            shape.line.color.rgb = COLOR["header_text"]
            shape.line.width = Pt(0.75)

            # Text formatting
            text_frame = shape.text_frame
            text_frame.text = full_title
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(14)
            run.font.name = FONT["body"]
            run.font.color.rgb = COLOR["text"]

            # Match by song number
            if song_number in number_index_map:
                song_index = number_index_map[song_number]
                try:
                    shape.click_action.target_slide = song_slide_map[song_index]
                except Exception as e:
                    print(f"Error linking '{full_title}' in alpha index: {e}")
            else:
                print(f"Warning: No match for song number {song_number} ({title})")

        return slide
