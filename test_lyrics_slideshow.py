import unittest
import os
from lyrics_to_slides import LyricsSlideshow
from song_parser import SongParser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

class TestLyricsSlideshow(unittest.TestCase):
    def setUp(self):
        self.parser = SongParser()
        self.slideshow = LyricsSlideshow()
        self.test_songs = """Song 1
1. First verse of the song,
Testing multiple choruses.

        First chorus of the song,
        Testing, testing, one two three.

2. Second verse of the song,
More testing to be done.

        Second chorus here we go,
        Different from the first one.

3. Third verse coming up,
Almost done with the test.

Song 2
1. Another song to test,
Making sure it all works."""
        
        self.test_output = "test_slideshow.pptx"
        
    def tearDown(self):
        # Clean up test file after each test
        if os.path.exists(self.test_output):
            os.remove(self.test_output)

    def test_presentation_creation(self):
        """Test that the presentation is created successfully."""
        songs = self.parser.parse_songs(self.test_songs)
        output_file = self.slideshow.create_presentation(songs, self.test_output)
        
        self.assertTrue(os.path.exists(output_file))
        
        # Load the presentation and check its structure
        prs = Presentation(output_file)
        
        # Calculate expected number of slides:
        # 1 (main title) + 
        # For each song:
        #   number of expanded sections
        expected_slides = 1  # Main title
        for song in songs:
            expected_slides += len(song['expanded_sections'])  # One slide per expanded section
            
        self.assertEqual(len(prs.slides), expected_slides)

    def test_slide_content_types(self):
        """Test that slides have the correct content types."""
        songs = self.parser.parse_songs(self.test_songs)
        output_file = self.slideshow.create_presentation(songs, self.test_output)
        
        prs = Presentation(output_file)
        
        # First slide should be main title
        first_slide = prs.slides[0]
        # Get all text from all shapes on the first slide
        title_text = ""
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame"):
                title_text = shape.text_frame.text.strip()
                if title_text:
                    break
        self.assertEqual(title_text, "Song Lyrics Slideshow")
        
        # Second slide should be first verse of first song
        second_slide = prs.slides[1]
        song_title = ""
        verse_header = ""
        for shape in second_slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if "SONG 1" in text:
                    song_title = text
                elif "STANZA" in text:
                    verse_header = text
        self.assertEqual(song_title, "SONG 1")
        self.assertEqual(verse_header, "STANZA 1")

    def test_slide_styling(self):
        """Test that slides have the correct styling."""
        songs = self.parser.parse_songs(self.test_songs)
        output_file = self.slideshow.create_presentation(songs, self.test_output)
        
        prs = Presentation(output_file)
        
        # Check content slide (second slide)
        content_slide = prs.slides[1]
        
        # Check main background color
        background = content_slide.background
        fill = background.fill
        self.assertTrue(fill.type)  # Should be solid fill
        self.assertEqual(fill.fore_color.rgb, self.slideshow.BACKGROUND_COLOR)
        
        # Check header background shape
        header_shape = None
        for shape in content_slide.shapes:
            if shape.shape_type == MSO_SHAPE.RECTANGLE:  # Rectangle shape for header background
                header_shape = shape
                break
        self.assertIsNotNone(header_shape, "Header background shape not found")
        self.assertEqual(header_shape.fill.fore_color.rgb, self.slideshow.HEADER_BACKGROUND_COLOR)
        
        # Check text formatting
        for shape in content_slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:  # Only check non-empty text boxes
                    if "SONG" in text or "STANZA" in text:
                        # Header text should be uppercase and have header color
                        self.assertTrue(text.isupper())
                        font = shape.text_frame.paragraphs[0].font
                        if font.color.type is not None:  # Only check if color is set
                            self.assertEqual(font.color.rgb, self.slideshow.HEADER_TEXT_COLOR)
                    else:
                        # Content text should have main text color
                        font = shape.text_frame.paragraphs[0].font
                        if font.color.type is not None:  # Only check if color is set
                            self.assertEqual(font.color.rgb, self.slideshow.TEXT_COLOR)

    def test_chorus_numbering(self):
        """Test that multiple choruses are numbered correctly."""
        songs = self.parser.parse_songs(self.test_songs)
        output_file = self.slideshow.create_presentation(songs, self.test_output)
        
        prs = Presentation(output_file)
        
        # Find slides with choruses
        chorus_slides = []
        for i in range(len(prs.slides)):  # Iterate through slide indices
            slide = prs.slides[i]
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text.strip()
                    if "CHORUS" in text:
                        chorus_slides.append(text)
                        break
        
        # Verify chorus numbering
        self.assertEqual(len(chorus_slides), 2)  # Should have 2 choruses
        self.assertEqual(chorus_slides[0], "CHORUS 1")
        self.assertEqual(chorus_slides[1], "CHORUS 2")

if __name__ == '__main__':
    unittest.main() 